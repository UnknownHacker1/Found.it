"""
Document Parser - Extracts text content from various file types
Supports: PDF, DOCX, PPTX, XLSX, TXT, and 60+ code/config formats
"""

from pathlib import Path
from typing import Optional
import logging

# PDF parsing
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

# Word document parsing
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

# PowerPoint parsing
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Excel parsing
try:
    import openpyxl
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class DocumentParser:
    """Extracts text content from various document types"""

    # Text-based file extensions that can be read directly
    TEXT_EXTENSIONS = {
        # Documents & Notes
        '.txt', '.md', '.markdown', '.log', '.rtf',

        # Data & Config
        '.csv', '.tsv', '.json', '.xml', '.yaml', '.yml', '.toml', '.ini', '.cfg', '.conf',

        # Code - C/C++
        '.c', '.cpp', '.cc', '.cxx', '.h', '.hpp', '.hxx',

        # Code - Web
        '.html', '.htm', '.css', '.scss', '.sass', '.less',
        '.js', '.jsx', '.ts', '.tsx', '.vue', '.svelte',

        # Code - Backend
        '.py', '.java', '.cs', '.go', '.rs', '.rb', '.php',

        # Code - Scripts
        '.sh', '.bash', '.zsh', '.bat', '.cmd', '.ps1',

        # Code - Mobile
        '.swift', '.kt', '.kts', '.m', '.mm',

        # Code - Other
        '.sql', '.r', '.scala', '.clj', '.ex', '.exs',

        # Markup & Data
        '.tex', '.bib', '.svg', '.graphql', '.proto',

        # Environment & Build
        '.env', '.gitignore', '.dockerignore', 'Dockerfile', 'Makefile',
        '.gradle', '.maven', '.cmake',

        # Documentation
        '.rst', '.adoc', '.textile'
    }

    def __init__(self):
        self.stats = {
            'pdf_support': PDF_AVAILABLE,
            'docx_support': DOCX_AVAILABLE,
            'pptx_support': PPTX_AVAILABLE,
            'xlsx_support': XLSX_AVAILABLE
        }

        # Log which parsers are available
        logger.info("Document Parser initialized:")
        logger.info(f"  PDF support: {'✓' if PDF_AVAILABLE else '✗ (install PyPDF2)'}")
        logger.info(f"  DOCX support: {'✓' if DOCX_AVAILABLE else '✗ (install python-docx)'}")
        logger.info(f"  PPTX support: {'✓' if PPTX_AVAILABLE else '✗ (install python-pptx)'}")
        logger.info(f"  XLSX support: {'✓' if XLSX_AVAILABLE else '✗ (install openpyxl)'}")

    def can_parse(self, file_path: Path) -> bool:
        """Check if the file type is supported"""
        ext = file_path.suffix.lower()

        if ext in self.TEXT_EXTENSIONS:
            return True
        if ext == '.pdf' and PDF_AVAILABLE:
            return True
        if ext == '.docx' and DOCX_AVAILABLE:
            return True
        if ext == '.pptx' and PPTX_AVAILABLE:
            return True
        if ext == '.xlsx' and XLSX_AVAILABLE:
            return True

        return False

    def parse(self, file_path: Path) -> Optional[str]:
        """
        Extract text content from a file
        Returns the extracted text or None if parsing fails
        """
        try:
            ext = file_path.suffix.lower()

            if ext in self.TEXT_EXTENSIONS:
                return self._parse_text_file(file_path)
            elif ext == '.pdf':
                return self._parse_pdf(file_path)
            elif ext == '.docx':
                return self._parse_docx(file_path)
            elif ext == '.pptx':
                return self._parse_pptx(file_path)
            elif ext == '.xlsx':
                return self._parse_xlsx(file_path)
            else:
                logger.warning(f"Unsupported file type: {ext}")
                return None

        except Exception as e:
            logger.error(f"Error parsing {file_path}: {e}")
            return None

    def _parse_text_file(self, file_path: Path) -> Optional[str]:
        """Read plain text files"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']

            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                        # Limit size to prevent memory issues
                        if len(content) > 1_000_000:  # 1MB of text
                            content = content[:1_000_000]
                        return content
                except UnicodeDecodeError:
                    continue

            logger.warning(f"Could not decode text file: {file_path}")
            return None

        except Exception as e:
            logger.error(f"Error reading text file {file_path}: {e}")
            return None

    def _parse_pdf(self, file_path: Path) -> Optional[str]:
        """Extract text from PDF files"""
        if not PDF_AVAILABLE:
            logger.warning("PyPDF2 not installed. Cannot parse PDF files.")
            return None

        try:
            text_content = []

            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)

                # Limit to first 50 pages to prevent excessive processing
                max_pages = min(len(pdf_reader.pages), 50)

                for page_num in range(max_pages):
                    page = pdf_reader.pages[page_num]
                    text = page.extract_text()
                    if text:
                        text_content.append(text)

            result = '\n'.join(text_content)
            return result if result.strip() else None

        except Exception as e:
            logger.error(f"Error parsing PDF {file_path}: {e}")
            return None

    def _parse_docx(self, file_path: Path) -> Optional[str]:
        """Extract text from Word documents"""
        if not DOCX_AVAILABLE:
            logger.warning("python-docx not installed. Cannot parse DOCX files.")
            return None

        try:
            doc = Document(file_path)
            text_content = []

            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)

            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text_content.append(cell.text)

            result = '\n'.join(text_content)
            return result if result.strip() else None

        except Exception as e:
            logger.error(f"Error parsing DOCX {file_path}: {e}")
            return None

    def _parse_pptx(self, file_path: Path) -> Optional[str]:
        """Extract text from PowerPoint presentations"""
        if not PPTX_AVAILABLE:
            logger.warning("python-pptx not installed. Cannot parse PPTX files.")
            return None

        try:
            prs = Presentation(file_path)
            text_content = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text)

            result = '\n'.join(text_content)
            return result if result.strip() else None

        except Exception as e:
            logger.error(f"Error parsing PPTX {file_path}: {e}")
            return None

    def _parse_xlsx(self, file_path: Path) -> Optional[str]:
        """Extract text from Excel spreadsheets"""
        if not XLSX_AVAILABLE:
            logger.warning("openpyxl not installed. Cannot parse XLSX files.")
            return None

        try:
            workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            text_content = []

            # Limit to first 10 sheets to prevent excessive processing
            max_sheets = min(len(workbook.sheetnames), 10)

            for sheet_name in workbook.sheetnames[:max_sheets]:
                sheet = workbook[sheet_name]

                # Add sheet name as header
                text_content.append(f"Sheet: {sheet_name}")

                # Limit to first 1000 rows to prevent excessive processing
                max_rows = min(sheet.max_row, 1000) if sheet.max_row else 1000

                for row in sheet.iter_rows(max_row=max_rows, values_only=True):
                    # Convert row values to strings and filter out None/empty
                    row_text = ' | '.join(str(cell) for cell in row if cell is not None and str(cell).strip())
                    if row_text:
                        text_content.append(row_text)

            workbook.close()

            result = '\n'.join(text_content)
            return result if result.strip() else None

        except Exception as e:
            logger.error(f"Error parsing XLSX {file_path}: {e}")
            return None

    def get_preview(self, content: str, max_length: int = 200) -> str:
        """Get a preview of the content"""
        if not content:
            return ""

        # Clean up whitespace
        preview = ' '.join(content.split())

        if len(preview) > max_length:
            preview = preview[:max_length] + "..."

        return preview
